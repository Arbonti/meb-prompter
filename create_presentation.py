import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    # 16:9 format
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Renk paleti
    NAVY = RGBColor(15, 23, 42) # #0f172a
    RED = RGBColor(192, 57, 43) # #c0392b
    LIGHT_BG = RGBColor(248, 250, 252) # #f8fafc
    WHITE = RGBColor(255, 255, 255)
    DARK_TEXT = RGBColor(30, 41, 59) # #1e293b
    MUTED_TEXT = RGBColor(100, 116, 139) # #64748b
    
    # 1. SLAYT: KAPAK (Koyu Arka Plan)
    slide_layout = prs.slide_layouts[6] # Boş slayt
    slide = prs.slides.add_slide(slide_layout)
    
    # Arka plan rengi
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY
    
    # MEB Logo Ekleme
    if os.path.exists("meb_logo.png"):
        try:
            # 13.333 inç genişliğin ortasına yerleştirme (genişlik 2.0 inç, sol = (13.333-2)/2 = 5.666)
            slide.shapes.add_picture("meb_logo.png", Inches(5.66), Inches(0.5), width=Inches(2.0))
        except Exception as e:
            print(f"Logo eklenirken hata oluştu: {e}")
            
    # Başlık Kutusu
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.333), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "AYDIN İL MİLLİ EĞİTİM MÜDÜRLÜĞÜ"
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = WHITE
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "Staj Çalışmaları ve Yapay Zeka Destekli İstatistik Bilgi Sistemi Projesi"
    p2.font.size = Pt(22)
    p2.font.color.rgb = RED
    p2.font.name = "Arial"
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(15)
    
    p3 = tf.add_paragraph()
    p3.text = "Hazırlayan: Nedim ALPTEKİN\nÖğrenci No: 2212903054\nTarih: Haziran 2026"
    p3.font.size = Pt(16)
    p3.font.color.rgb = MUTED_TEXT
    p3.font.name = "Arial"
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(30)
    
    # Diğer slaytlar için genel fonksiyon
    def add_content_slide(title, bullets, image_path=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Arka plan
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = LIGHT_BG
        
        # Başlık
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(1.0))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title
        p_title.font.bold = True
        p_title.font.size = Pt(28)
        p_title.font.color.rgb = NAVY
        p_title.font.name = "Arial"
        
        # Kırmızı alt çizgi
        line = slide.shapes.add_shape(
            1, # MSO_SHAPE.RECTANGLE (1 is rectangle)
            Inches(0.75), Inches(1.3), Inches(2.0), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RED
        line.line.color.rgb = RED
        
        # Sütunları Ayarla
        if image_path and os.path.exists(image_path):
            # İki Sütunlu Düzen (Sol: Metin, Sağ: Resim)
            content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(6.5), Inches(5.0))
            try:
                # Resmi yerleştir (Sağ sütun)
                # Genişlik 5.0 inç, sol = 7.5 inç, üst = 1.8 inç
                slide.shapes.add_picture(image_path, Inches(7.5), Inches(1.8), width=Inches(5.0))
            except Exception as e:
                print(f"Resim eklenirken hata oluştu ({image_path}): {e}")
        else:
            # Tek Sütunlu Düzen (Tam Genişlik)
            content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5.0))
            
        tf_content = content_box.text_frame
        tf_content.word_wrap = True
        
        for idx, bullet in enumerate(bullets):
            p = tf_content.add_paragraph() if idx > 0 else tf_content.paragraphs[0]
            p.text = bullet
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_TEXT
            p.font.name = "Arial"
            p.space_after = Pt(14)
            p.level = 0
            
            # Eğer alt başlık varsa (girintili)
            if bullet.startswith("  -") or bullet.startswith("    "):
                p.text = bullet.replace("  -", "").replace("    ", "").strip()
                p.level = 1
                p.font.size = Pt(16)
                p.font.color.rgb = MUTED_TEXT
                
    # 2. SLAYT
    add_content_slide(
        "Staj Yapılan Kurum ve Birim",
        [
            "• Aydın İl Milli Eğitim Müdürlüğü (Aydın MEM)",
            "  - İl genelindeki 17 ilçede bulunan okul, öğrenci, öğretmen ve eğitim tesislerinin yönetiminden sorumludur.",
            "• Bilgi İşlem ve İstatistik Şubesi",
            "  - Kurumun dijital altyapısını yönetir.",
            "  - e-Okul, MEBBİS ve diğer resmi veri tabanlarındaki istatistiksel verilerin doğruluğunu, saklanmasını ve raporlanmasını koordine eder.",
            "• Stajın Temel Amacı",
            "  - Eğitim istatistiklerinin sorgulanması ve raporlanması süreçlerini modernize etmek, kullanıcı dostu hale getirmek."
        ]
    )
    
    # 3. SLAYT
    add_content_slide(
        "Mevcut Sorunlar ve Proje İhtiyacı",
        [
            "• Klasik Veri Raporlama Zorlukları",
            "  - e-Okul ve MEBBİS verilerinin genellikle statik tablolar halinde olması ve arama yapmanın uzmanlık gerektirmesi.",
            "  - Karar vericilerin (Müdür, Şube Müdürü vb.) anlık verilere (örn: 'ilçedeki yabancı uyruklu öğrenci sayısı') ulaşmakta zaman kaybetmesi.",
            "• Çözüm Odaklı Yaklaşım: MEB Prompter",
            "  - Doğal dil işleme (NLP) yeteneğine sahip, arama motoru kolaylığında çalışan bir web uygulaması.",
            "  - Kullanıcıların sadece Türkçe soru sorarak (örn: 'Kuşadası'nda kaç ilkokul var?') istatistiklere anında erişebilmesi."
        ]
    )
    
    # 4. SLAYT (Görselli Slayt)
    add_content_slide(
        "Geliştirilen Sistem: MEB İstatistik Asistanı",
        [
            "• Yapay Zeka Entegrasyonu (Gemini AI)",
            "  - Google Gemini 2.0 Flash modeli kullanılarak doğal dildeki sorular analiz edilir ve yanıt üretilir.",
            "• İnteraktif Yoğunluk Haritası",
            "  - Leaflet.js kütüphanesi kullanılarak öğretmen, öğrenci ve okul sayıları ilçe bazında görselleştirilir.",
            "• Akıllı Filtreleme Sistemi (Okul Sorgula)",
            "  - Deprem risk analizi, derslik, taşımalı eğitim gibi 30'dan fazla kriterle analiz."
        ],
        image_path="images/media__1780309827685.png"
    )
    
    # 5. SLAYT
    add_content_slide(
        "Teknik Altyapı ve Mimarisi",
        [
            "• Ön Yüz (Frontend): HTML5, CSS3, Vanilla JS",
            "  - Kurumsal ağ kısıtlamalarına takılmaması için tamamen istemci tarafında (Client-Side) çalışan hafif mimari.",
            "• Yapay Zeka ve Veri Bağlamı (Prompt Context):",
            "  - Kullanıcının sorusuna göre MEB veritabanından ilgili özet metin oluşturulur ve Gemini API'ye bağlam olarak beslenir.",
            "• Harita Entegrasyonları:",
            "  - Leaflet.js ile coğrafi veri haritalandırma tecrübesi."
        ]
    )
    
    # 6. SLAYT (Görselli Slayt)
    add_content_slide(
        "Staj Süresince Yaptığım Çalışmalar & Katkılarım",
        [
            "• 1. Dinamik Belge ve Veri Yükleme Modülü",
            "  - Excel (.xlsx), Word (.docx) ve PDF (.pdf) belgelerinin sürükle-bırak ile sisteme yüklenmesi sağlandı.",
            "• 2. Canlı Veri Senkronizasyonu",
            "  - Excel'den yüklenen yeni verilerin sorgu tablosunda, istatistik kartlarında ve haritada canlı güncellenmesi sağlandı.",
            "• 3. Belge Analiz Sistemi (RAG)",
            "  - PDF ve Word dosyalarındaki metinlerin ayıklanarak Gemini API sohbet bağlamına eklenmesi sağlandı (25.000 karakterlik limit)."
        ],
        image_path="images/media__1780310518137.png"
    )
    
    # 7. SLAYT
    add_content_slide(
        "Edinilen Deneyimler ve Kazanımlar",
        [
            "• Teknik Kazanımlar:",
            "  - Google Gemini AI API ile entegrasyon deneyimi ve Prompt mühendisliği.",
            "  - SheetJS (XLSX), Mammoth.js ve PDF.js ile istemci tarafında dosya okuma ve veri işleme yetkinliği.",
            "  - Leaflet.js ile coğrafi veri haritalandırma tecrübesi.",
            "• Profesyonel Kazanımlar:",
            "  - Kurumsal ihtiyaç analizi yapma ve kullanıcı deneyimine (UX) göre arayüz geliştirme.",
            "  - JavaScript tabanlı asenkron programlama ve tarayıcı hafıza (localStorage) yönetimi."
        ]
    )
    
    # 8. SLAYT: KAPANIŞ (Koyu Arka Plan)
    slide = prs.slides.add_slide(slide_layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = NAVY
    
    # MEB Logo Ekleme
    if os.path.exists("meb_logo.png"):
        try:
            slide.shapes.add_picture("meb_logo.png", Inches(5.66), Inches(0.5), width=Inches(2.0))
        except Exception as e:
            pass
            
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.333), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "TEŞEKKÜRLER"
    p.font.bold = True
    p.font.size = Pt(44)
    p.font.color.rgb = WHITE
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "Sorularınız ve Değerlendirmeleriniz"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RED
    p2.font.name = "Arial"
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(15)
    
    p3 = tf.add_paragraph()
    p3.text = "Nedim ALPTEKİN\nÖğrenci No: 2212903054"
    p3.font.size = Pt(16)
    p3.font.color.rgb = MUTED_TEXT
    p3.font.name = "Arial"
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(30)
    
    # Sunumları kaydet
    prs.save("Staj_Sunumu.pptx")
    print("[Success] Presentation updated with images and author details!")

if __name__ == "__main__":
    create_presentation()
